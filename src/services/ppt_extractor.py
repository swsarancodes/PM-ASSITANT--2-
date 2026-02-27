from pptx import Presentation
from io import BytesIO
import re
import unicodedata


def clean_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r'[\x00-\x1F\x7F]', '', text)
    return text


def extract_text_from_ppt(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))

    slides_text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)

    combined_text = "\n".join(slides_text)

    return clean_text(combined_text)